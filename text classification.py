#!/usr/bin/env python
# coding: utf-8

# In[46]:


from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.cluster import KMeans
import os
import nltk
from sklearn.metrics import silhouette_score
nltk.download('stopwords')
from nltk.stem.porter import PorterStemmer
from nltk.corpus import stopwords
ps = PorterStemmer()
import numpy as np
import pandas as pd
import glob
import jieba
import re
import csv
import string
import nltk
#stop = set(stopwords.words('english'))
#提供自己定義的stopwords
file1 = open('C:\\Users\\user\\Desktop\\stopwords.csv','r')
stopword = csv.reader(file1,delimiter = ';')
stopwords = []
for i in stopword:
    stopwords.append(i)
topic = []
texts = []
#搜括所有的ppt檔
for root,dirs, files in os.walk('.'):
    for name in files:
        if os.path.splitext(name)[1] == '.pptx':
            topic.append(os.path.join(root,name))
            prs = Presentation(os.path.join(root,name))
            text = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            texts.append(text)
"""
for eachfile in glob.glob('C:\\Users\\user\\Desktop\\*.pptx'):
    prs = Presentation(eachfile)
    topic.append(eachfile)
    text = []
    #print("----------------------")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    texts.append(text)
"""
#文字前處理, 去除 stopword, 中文字串分割
textss = []
for i in texts:
    i = ' '.join(i)
    i = re.sub('\d','',i)
    i = re.sub("[\s+\.\!\/_,$%^*(+\"\']+|[+【】:)——！，。？?、~@#￥%……&*（）：]+", " ", i)
    i = i.lower()
    i = ' '.join([ps.stem(w) for w in i.split() if w not in stopwords])
    i = jieba.cut(i,cut_all = False)
    textss.append(" ".join(i))
 
vectorizer = TfidfVectorizer(stop_words='english')
#轉換 tf-idf 矩陣
X = vectorizer.fit_transform(textss)
X_array = X.toarray()
X_name = vectorizer.get_feature_names()
tf_idf_matrix = pd.DataFrame(X_array,columns = X_name)
tf_idf_matrix.to_csv('C:\\Users\\user\\Desktop\\TF_IDF.csv',index = None,encoding = 'utf-8')
#Silhouette 
range_n_clusters = list (range(2,10))
scores = []
for n_clusters in range_n_clusters:
    clusterer = KMeans(n_clusters=n_clusters)
    preds = clusterer.fit_predict(X)
    centers = clusterer.cluster_centers_
    score = silhouette_score(X, preds)
    scores.append(score)

scores = pd.DataFrame(scores, columns = {'score'},index = range_n_clusters)
print (scores)
scores.transpose().plot(kind=  'bar')
silhouette_order = scores.sort_values(by= 'score', ascending = False)
true_k = silhouette_order.index.values.tolist()[0]

model = KMeans(n_clusters=true_k, init='k-means++', max_iter=300, n_init=1)
model.fit(X)
cluster = pd.DataFrame(model.labels_)

texts = pd.DataFrame(textss, columns = {'content'})
topic = pd.DataFrame(topic, columns = {'file'})
texts = pd.concat([topic,texts],axis = 1)
texts['cluster_group'] = cluster
order_centroids = model.cluster_centers_.argsort()[:, ::-1]
#print (order_centroids)
terms = vectorizer.get_feature_names()
for a in range(0, len(texts)):
    for i in range(true_k):
        key = []
        for ind in order_centroids[i, :10]:
            key.append(terms[ind])
        if texts.iloc[a]['cluster_group'] == i:
            texts.loc[a,'keywords'] = str(' '.join(key))
texts = texts.sort_values(by = 'cluster_group',ascending = True) 
texts.to_excel('C:\\Users\\user\\Desktop\\result.xlsx',index = None,encoding = 'utf-8')


# In[ ]:




